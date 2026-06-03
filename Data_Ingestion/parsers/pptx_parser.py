"""
PowerPoint / PPTX Parser — python-pptx
========================================
Extracts per slide: title, body text, images, tables, and speaker notes.

Production considerations:
  - Grouped shapes are traversed recursively — architecture diagrams and flowcharts
    in PowerPoint are almost always grouped shapes; without this they are silently skipped
  - Speaker notes are extracted as a body text element (often contain key context)
  - Image alt-text (description attribute) is used as caption before falling back to shape name
  - SmartArt text is extracted from the XML text frame even when shape_type != PICTURE
  - Per-element try/except — a bad shape never aborts the slide or deck
"""

from __future__ import annotations
import base64
import logging
from datetime import datetime
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from models.meta_schema import (
    ParsedDocument, TextElement, ImageElement, TableElement, SourceLocation,
)
from parsers.base_parser import BaseParser

logger = logging.getLogger(__name__)


class PPTXParser(BaseParser):

    def parse(self, file_path: str | Path) -> ParsedDocument:
        file_path = Path(file_path)
        doc_meta  = self._base_doc(file_path, "pptx")

        try:
            prs = Presentation(str(file_path))
        except Exception as e:
            raise RuntimeError(f"Cannot open PPTX '{file_path.name}': {e}") from e

        doc_meta.summary.slide_count = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            order = 0

            # ── Slide title ───────────────────────────────────────────────────
            try:
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    title_text = self._clean(slide.shapes.title.text)
                    if title_text:
                        elem = TextElement.build(
                            idx           = self._next_text_idx(),
                            content       = title_text,
                            loc           = SourceLocation(slide=slide_num, page=slide_num, order=order),
                            heading_level = 1,
                            is_title      = True,
                        )
                        doc_meta.text_elements.append(elem)
                        order += 1
            except Exception as e:
                logger.warning("Slide %d title failed: %s", slide_num, e)

            # ── All shapes (recursive — handles groups) ───────────────────────
            for shape in slide.shapes:
                order = self._process_shape(shape, slide_num, doc_meta, order, is_title_shape=shape == slide.shapes.title)

            # ── Speaker notes ─────────────────────────────────────────────────
            try:
                if slide.has_notes_slide:
                    notes_tf = slide.notes_slide.notes_text_frame
                    notes_raw = self._clean(notes_tf.text) if notes_tf else ""
                    if notes_raw:
                        elem = TextElement.build(
                            idx       = self._next_text_idx(),
                            content   = f"[Speaker Notes] {notes_raw}",
                            loc       = SourceLocation(slide=slide_num, page=slide_num, order=order),
                        )
                        doc_meta.text_elements.append(elem)
                        order += 1
            except Exception as e:
                logger.warning("Slide %d notes failed: %s", slide_num, e)

        doc_meta.parsed_at = datetime.utcnow()
        doc_meta.rebuild_summary()
        return doc_meta

    # ── Shape processor (called recursively for groups) ───────────────────────

    def _process_shape(
        self,
        shape,
        slide_num: int,
        doc_meta: ParsedDocument,
        order:    int,
        is_title_shape: bool = False,
    ) -> int:
        """
        Process a single shape. Returns the updated order counter.
        Recurses into GROUP shapes so nested content is never missed.
        """
        try:
            # ── GROUP — recurse into children ─────────────────────────────────
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    order = self._process_shape(child, slide_num, doc_meta, order)
                return order

            loc = SourceLocation(slide=slide_num, page=slide_num, order=order)

            # ── Picture ───────────────────────────────────────────────────────
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_blob = shape.image.blob
                    ext      = (shape.image.ext or "png").lower()
                    b64      = base64.b64encode(img_blob).decode("utf-8")
                    caption  = _shape_alt_text(shape) or None

                    elem = ImageElement.build(
                        idx         = self._next_img_idx(),
                        loc         = loc,
                        base64_data = b64,
                        format      = ext,
                        size_bytes  = len(img_blob),
                        caption     = caption,
                    )
                    doc_meta.image_elements.append(elem)
                    order += 1
                except Exception as e:
                    logger.warning("Slide %d picture failed: %s", slide_num, e)
                return order

            # ── Table ─────────────────────────────────────────────────────────
            if shape.has_table:
                try:
                    tbl       = shape.table
                    rows_data = []
                    for row in tbl.rows:
                        rows_data.append([
                            self._clean(cell.text) for cell in row.cells
                        ])

                    if rows_data:
                        headers = rows_data[0]
                        rows    = rows_data[1:]
                        caption = _shape_alt_text(shape) or self._clean(shape.name) or None

                        elem = TableElement.build(
                            idx     = self._next_tbl_idx(),
                            loc     = loc,
                            headers = headers,
                            rows    = rows,
                            caption = caption,
                        )
                        doc_meta.table_elements.append(elem)
                        order += 1
                except Exception as e:
                    logger.warning("Slide %d table failed: %s", slide_num, e)
                return order

            # ── Text frame (non-title) ─────────────────────────────────────────
            # Covers regular text boxes AND SmartArt (which has a text frame even
            # though its visual representation is handled differently)
            if shape.has_text_frame and not is_title_shape:
                try:
                    for para in shape.text_frame.paragraphs:
                        raw = self._clean(para.text)
                        if not raw:
                            continue
                        lvl       = para.level   # indent level (0 = top)
                        h         = 2 if lvl == 0 else None
                        is_bullet = lvl > 0

                        elem = TextElement.build(
                            idx           = self._next_text_idx(),
                            content       = raw,
                            loc           = SourceLocation(slide=slide_num, page=slide_num, order=order),
                            heading_level = h,
                            is_bullet     = is_bullet,
                        )
                        doc_meta.text_elements.append(elem)
                        order += 1
                except Exception as e:
                    logger.warning("Slide %d text frame failed: %s", slide_num, e)

        except Exception as e:
            logger.warning("Slide %d shape processing failed: %s", slide_num, e)

        return order


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def _shape_alt_text(shape) -> Optional[str]:
    """
    Return the alt-text / description of a shape if set by the author.
    Alt text is stored in the cNvPr element's 'descr' attribute.
    Returns None if not set or empty.
    """
    try:
        # python-pptx exposes this via shape._element for any shape type
        sp_el = shape._element
        # Try common cNvPr locations (pic:nvPicPr/pic:cNvPr, p:nvSpPr/p:cNvPr, etc.)
        for cNvPr in sp_el.iter():
            local = cNvPr.tag.split("}")[-1] if "}" in cNvPr.tag else cNvPr.tag
            if local == "cNvPr":
                desc = cNvPr.get("descr", "").strip()
                if desc:
                    return desc
                break
    except Exception:
        pass
    return None
